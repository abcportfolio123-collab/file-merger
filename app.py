import os
import subprocess
import tempfile
import shutil
import uuid
from flask import Flask, request, render_template, send_file, jsonify

from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from PIL import Image
from reportlab.pdfgen import canvas as rl_canvas

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB safety cap, not a "real" limit

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "uploads")
OUTPUT_DIR = os.path.join(BASE_DIR, "outputs")
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

OFFICE_EXTS = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".odt", ".odp", ".ods", ".rtf"}
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp", ".gif"}
PDF_EXT = ".pdf"


def convert_office_to_pdf(input_path, out_dir):
    """Use LibreOffice headless to convert any office doc to PDF, preserving formatting."""
    cmd = [
        "soffice", "--headless", "--norestore",
        "--convert-to", "pdf", "--outdir", out_dir, input_path
    ]
    subprocess.run(cmd, check=True, timeout=180,
                    stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    base = os.path.splitext(os.path.basename(input_path))[0]
    produced = os.path.join(out_dir, base + ".pdf")
    if not os.path.exists(produced):
        raise RuntimeError(f"Conversion failed for {input_path}")
    return produced


def image_to_pdf(input_path, out_dir):
    """Convert an image to a single-page PDF sized to the image (fit to a normal page)."""
    img = Image.open(input_path).convert("RGB")
    base = os.path.splitext(os.path.basename(input_path))[0]
    out_path = os.path.join(out_dir, base + "_img.pdf")

    page_w, page_h = letter
    img_w, img_h = img.size
    scale = min((page_w - 72) / img_w, (page_h - 72) / img_h, 1.0) if img_w > 0 and img_h > 0 else 1.0
    draw_w, draw_h = img_w * scale, img_h * scale
    x = (page_w - draw_w) / 2
    y = (page_h - draw_h) / 2

    tmp_img = os.path.join(out_dir, base + "_tmp.jpg")
    img.save(tmp_img, "JPEG", quality=95)

    c = rl_canvas.Canvas(out_path, pagesize=letter)
    c.drawImage(tmp_img, x, y, width=draw_w, height=draw_h)
    c.save()
    os.remove(tmp_img)
    return out_path


def build_notes_pdf(pptx_path, out_dir):
    """
    Build one selectable-text PDF page per slide containing that slide's speaker notes.
    Returns path to the notes PDF, or None if the pptx has no notes at all.
    """
    prs = Presentation(pptx_path)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle('NotesTitle', parent=styles['Heading2'], spaceAfter=12)
    body_style = ParagraphStyle('NotesBody', parent=styles['Normal'], fontSize=11, leading=15)

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    out_path = os.path.join(out_dir, base + "_notes.pdf")

    doc = SimpleDocTemplate(out_path, pagesize=letter,
                             topMargin=54, bottomMargin=54, leftMargin=54, rightMargin=54)
    story = []
    any_notes = False

    for idx, slide in enumerate(prs.slides, start=1):
        notes_text = ""
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            notes_text = tf.text.strip() if tf else ""
        if notes_text:
            any_notes = True
        story.append(Paragraph(f"Speaker Notes — Slide {idx}", title_style))
        display_text = notes_text if notes_text else "(No speaker notes for this slide)"
        # Escape basic XML special chars for reportlab Paragraph
        display_text = (display_text.replace("&", "&amp;")
                         .replace("<", "&lt;").replace(">", "&gt;"))
        display_text = display_text.replace("\n", "<br/>")
        story.append(Paragraph(display_text, body_style))
        story.append(Spacer(1, 24))
        if idx != len(prs.slides):
            from reportlab.platypus import PageBreak
            story.append(PageBreak())

    if not any_notes:
        return None

    doc.build(story)
    return out_path


def interleave_slides_and_notes(slides_pdf_path, notes_pdf_path, out_dir, base_name):
    """Produce slide1, notes1, slide2, notes2, ... as one PDF."""
    slides_reader = PdfReader(slides_pdf_path)
    notes_reader = PdfReader(notes_pdf_path) if notes_pdf_path else None

    writer = PdfWriter()
    for i, page in enumerate(slides_reader.pages):
        writer.add_page(page)
        if notes_reader and i < len(notes_reader.pages):
            writer.add_page(notes_reader.pages[i])

    out_path = os.path.join(out_dir, base_name + "_with_notes.pdf")
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path


def process_single_file(file_path, work_dir):
    """Convert one uploaded file to a PDF path (or list of PDFs), returns list of pdf paths to append."""
    ext = os.path.splitext(file_path)[1].lower()

    if ext == PDF_EXT:
        return [file_path]

    if ext in IMAGE_EXTS:
        return [image_to_pdf(file_path, work_dir)]

    if ext in OFFICE_EXTS:
        slides_pdf = convert_office_to_pdf(file_path, work_dir)
        if ext == ".pptx" or ext == ".ppt":
            try:
                notes_pdf = build_notes_pdf(file_path, work_dir) if ext == ".pptx" else None
            except Exception:
                notes_pdf = None
            if notes_pdf:
                base_name = os.path.splitext(os.path.basename(file_path))[0]
                combined = interleave_slides_and_notes(slides_pdf, notes_pdf, work_dir, base_name)
                return [combined]
        return [slides_pdf]

    raise ValueError(f"Unsupported file type: {ext}")


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/merge", methods=["POST"])
def merge():
    files = request.files.getlist("files")
    if not files or len(files) == 0:
        return jsonify({"error": "No files uploaded"}), 400

    job_id = str(uuid.uuid4())
    work_dir = os.path.join(UPLOAD_DIR, job_id)
    os.makedirs(work_dir, exist_ok=True)

    saved_paths = []
    try:
        for f in files:
            if f.filename == "":
                continue
            safe_name = f.filename.replace("/", "_").replace("\\", "_")
            path = os.path.join(work_dir, safe_name)
            f.save(path)
            saved_paths.append(path)

        if not saved_paths:
            return jsonify({"error": "No valid files uploaded"}), 400

        writer = PdfWriter()
        for path in saved_paths:
            pdf_parts = process_single_file(path, work_dir)
            for part in pdf_parts:
                reader = PdfReader(part)
                for page in reader.pages:
                    writer.add_page(page)

        output_path = os.path.join(OUTPUT_DIR, f"merged_{job_id}.pdf")
        with open(output_path, "wb") as f:
            writer.write(f)

        return send_file(output_path, as_attachment=True, download_name="merged.pdf")

    except subprocess.CalledProcessError as e:
        return jsonify({"error": f"Conversion failed: {e.stderr.decode(errors='ignore')[:500]}"}), 500
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 7860))
    app.run(host="0.0.0.0", port=port)
